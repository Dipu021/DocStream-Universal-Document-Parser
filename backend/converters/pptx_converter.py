from pathlib import Path
from typing import AsyncIterator

from pptx import Presentation


async def convert(path: Path) -> AsyncIterator[dict]:
    try:
        prs = Presentation(str(path))
    except Exception as e:
        yield {
            "section": "Error",
            "markdown": f"> **Could not open this PowerPoint file** — it may be corrupt or not a valid .pptx.\n>\n> `{e}`\n",
            "progress": 1.0,
        }
        return

    try:
        slides = list(prs.slides)
    except Exception as e:
        yield {
            "section": "Error",
            "markdown": f"> **Could not read slides.**\n>\n> `{e}`\n",
            "progress": 1.0,
        }
        return

    total = max(len(slides), 1)

    for i, slide in enumerate(slides):
        lines = [f"## Slide {i + 1}"]

        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None

        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if not text:
                            continue
                        if shape is title_shape:
                            lines.append(f"### {text}")
                        else:
                            lines.append(f"- {text}")

                if shape.has_table:
                    tbl = shape.table
                    rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
                    if rows:
                        width = max(len(r) for r in rows)
                        rows = [r + [""] * (width - len(r)) for r in rows]
                        header = rows[0]
                        lines.append("| " + " | ".join(header) + " |")
                        lines.append("| " + " | ".join(["---"] * width) + " |")
                        for row in rows[1:]:
                            safe_row = [c.replace("|", "\\|").replace("\n", " ") for c in row]
                            lines.append("| " + " | ".join(safe_row) + " |")
            except Exception as e:
                lines.append(f"<!-- skipped one shape due to error: {e} -->")

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes = f"\n> Notes: {slide.notes_slide.notes_text_frame.text.strip()}"
        except Exception:
            pass

        md = "\n\n".join(lines) + notes + "\n"

        yield {
            "section": f"Slide {i + 1}",
            "markdown": md,
            "progress": (i + 1) / total,
        }